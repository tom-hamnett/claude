"""Multi-format intake — read and interpret ANYTHING the user brings.

The content intake shouldn't care what format the source is. A CV (PDF/Word), a
deck (PowerPoint), a spreadsheet, a photo, a video, a plain paste, or a link — all
of it should become usable reference text for the generator.

Routing:
  • tabular (CSV/TSV/XLSX)     → local table reader (numbers stay numbers)
  • text (TXT/MD)              → read directly
  • PDF / Word / PowerPoint    → local text extraction; PDF falls back to Gemini
                                 (scanned/complex) when there's no extractable text
  • images / video / audio     → Gemini multimodal (describe / transcribe / extract)
  • URLs                       → fetch + strip to text

Everything is best-effort and guarded: a missing library or key degrades to a note,
never an exception. Extracted text is stored in the Data Vault as one reference
source per batch, which the generator reads as context.
"""

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_PER_FILE = 8000        # cap each source's extracted text (keeps prompts bounded)
MAX_COMBINED = 24000       # cap the combined reference blob

TAB_EXT = {".csv", ".tsv", ".xlsx", ".xlsm", ".xltx"}
TXT_EXT = {".txt", ".md", ".markdown", ".rtf"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".heic"}
VID_EXT = {".mp4", ".mov", ".webm", ".m4v", ".mkv", ".avi"}
AUD_EXT = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac"}


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        return "\n".join((pg.extract_text() or "") for pg in reader.pages).strip()
    except Exception as e:
        logger.info("pdf read failed %s: %s", path.name, e)
        return ""


def _read_docx(path: Path) -> str:
    try:
        import docx
        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip()
    except Exception as e:
        logger.info("docx read failed %s: %s", path.name, e)
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"— Slide {i} —"]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            lines.append(t)
            if len(lines) > 1:
                out.append("\n".join(lines))
        return "\n\n".join(out).strip()
    except Exception as e:
        logger.info("pptx read failed %s: %s", path.name, e)
        return ""


def _read_text(path: Path) -> str:
    try:
        return path.read_text(errors="ignore").strip()
    except Exception:
        return ""


def interpret_upload(path: str, name: str = "") -> tuple[str, str]:
    """Read one file into (text, source_type). source_type ∈
    dataset|document|image|video|audio. '' text means nothing could be read."""
    p = Path(path)
    name = name or p.name
    ext = p.suffix.lower()
    text, stype = "", "document"
    try:
        if ext in TAB_EXT:
            from gtm_engine.video.data_insight import read_tabular, table_to_text
            cols, rows = read_tabular(str(p))
            text, stype = (table_to_text(cols, rows, max_rows=200), "dataset")
        elif ext in TXT_EXT:
            text = _read_text(p)
        elif ext in PDF_EXT:
            text = _read_pdf(p)
            if len(text) < 80:            # scanned / image-only PDF → let Gemini read it
                from gtm_engine.utils.media import interpret_file_gemini
                text = interpret_file_gemini(str(p)) or text
        elif ext in DOCX_EXT:
            text = _read_docx(p)
        elif ext in PPTX_EXT:
            text = _read_pptx(p)
        elif ext in IMG_EXT:
            from gtm_engine.utils.media import interpret_file_gemini
            text, stype = interpret_file_gemini(str(p)), "image"
        elif ext in VID_EXT:
            from gtm_engine.utils.media import interpret_file_gemini
            text, stype = interpret_file_gemini(str(p), is_video=True), "video"
        elif ext in AUD_EXT:
            from gtm_engine.utils.media import interpret_file_gemini
            text, stype = interpret_file_gemini(str(p)), "audio"
        else:
            # Unknown — try Gemini, then a plain read.
            from gtm_engine.utils.media import interpret_file_gemini
            text = interpret_file_gemini(str(p)) or _read_text(p)
    except Exception as e:
        logger.info("interpret_upload failed %s: %s", name, e)
        text = ""
    return (text or "").strip()[:MAX_PER_FILE], stype


def _html_to_text(html: str) -> str:
    html = re.sub(r"(?is)<(script|style|noscript|svg|head).*?</\1>", " ", html)
    html = re.sub(r"(?s)<[^>]+>", " ", html)
    import html as _h
    txt = _h.unescape(html)
    return re.sub(r"\s+", " ", txt).strip()


def interpret_url(url: str) -> str:
    """Fetch a link and return its readable text (best-effort)."""
    url = (url or "").strip()
    if not url:
        return ""
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    try:
        import httpx
        r = httpx.get(url, follow_redirects=True, timeout=25,
                      headers={"User-Agent": "Mozilla/5.0 (content-intake)"})
        if r.status_code != 200 or not r.text:
            return ""
        ctype = r.headers.get("content-type", "")
        text = r.text if "html" not in ctype else _html_to_text(r.text)
        return text.strip()[:MAX_PER_FILE]
    except Exception as e:
        logger.info("interpret_url failed %s: %s", url, e)
        return ""


def ingest_references(files: list[tuple] | None = None, links: list[str] | None = None,
                      name: str = "Reference material") -> tuple[int | None, list[str]]:
    """Interpret a mix of uploaded files [(path, name), ...] and links, combine into
    ONE Data Vault reference source, and return (source_id, notes). notes lists what
    was read (and anything that couldn't be). source_id is None if nothing usable."""
    files = files or []
    links = links or []
    sections, notes = [], []
    for path, fname in files:
        text, stype = interpret_upload(path, fname)
        if text:
            sections.append(f"### SOURCE: {fname} ({stype})\n{text}")
            notes.append(f"✓ read {fname} ({stype}, {len(text)} chars)")
        else:
            notes.append(f"⚠ couldn't read {fname} — add a Google key for images/video, "
                         "or paste the text")
    for url in links:
        text = interpret_url(url)
        if text:
            sections.append(f"### SOURCE (link): {url}\n{text}")
            notes.append(f"✓ read {url} ({len(text)} chars)")
        else:
            notes.append(f"⚠ couldn't fetch {url}")
    if not sections:
        return None, notes
    combined = "\n\n".join(sections)[:MAX_COMBINED]
    try:
        from gtm_engine.data_vault import DataVault, DataSource
        sid = DataVault().create(DataSource(
            name=name, description=f"Intake reference — {len(sections)} source(s)",
            source_type="document", content=combined, freshness="uploaded"))
        return sid, notes
    except Exception as e:
        logger.error("ingest_references store failed: %s", e)
        return None, notes
